#!/usr/bin/env python3
"""Create curated PPTX + manifest fixtures for adjustment-heavy shape benchmarks.

The first scenario focuses on curved-arrow presets because their adjusted geometry
is both visually sensitive and historically under-covered by our image checks.
Use the generated deck with evaluate/reference_render.py, evaluate/candidate_render.py,
and evaluate/shape_actual_coverage.py to compare renderer output against a rendered
reference image set.
"""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class ShapeVariant:
    label: str
    shape_type: MSO_AUTO_SHAPE_TYPE
    adjustments: dict[str, int]


CURVED_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "CURVED_RIGHT_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.CURVED_RIGHT_ARROW,
        {"adj1": 12_000, "adj2": 70_000, "adj3": 18_000},
    ),
    ShapeVariant(
        "CURVED_RIGHT_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.CURVED_RIGHT_ARROW,
        {"adj1": 42_000, "adj2": 30_000, "adj3": 42_000},
    ),
    ShapeVariant(
        "CURVED_LEFT_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.CURVED_LEFT_ARROW,
        {"adj1": 12_000, "adj2": 70_000, "adj3": 18_000},
    ),
    ShapeVariant(
        "CURVED_LEFT_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.CURVED_LEFT_ARROW,
        {"adj1": 42_000, "adj2": 30_000, "adj3": 42_000},
    ),
    ShapeVariant(
        "CURVED_UP_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.CURVED_UP_ARROW,
        {"adj1": 12_000, "adj2": 70_000, "adj3": 18_000},
    ),
    ShapeVariant(
        "CURVED_UP_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.CURVED_UP_ARROW,
        {"adj1": 42_000, "adj2": 30_000, "adj3": 42_000},
    ),
    ShapeVariant(
        "CURVED_DOWN_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.CURVED_DOWN_ARROW,
        {"adj1": 12_000, "adj2": 70_000, "adj3": 18_000},
    ),
    ShapeVariant(
        "CURVED_DOWN_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.CURVED_DOWN_ARROW,
        {"adj1": 42_000, "adj2": 30_000, "adj3": 42_000},
    ),
)

UP_DOWN_ARROW_CALLOUT_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "UP_DOWN_ARROW_CALLOUT_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW_CALLOUT,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000, "adj4": 15_000},
    ),
    ShapeVariant(
        "UP_DOWN_ARROW_CALLOUT_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW_CALLOUT,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 35_000, "adj4": 35_000},
    ),
    ShapeVariant(
        "UP_DOWN_ARROW_CALLOUT_ADJ_LONG",
        MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW_CALLOUT,
        {"adj1": 20_000, "adj2": 50_000, "adj3": 25_000, "adj4": 50_000},
    ),
    ShapeVariant(
        "UP_DOWN_ARROW_CALLOUT_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW_CALLOUT,
        {"adj1": 45_000, "adj2": 20_000, "adj3": 45_000, "adj4": 20_000},
    ),
)

WAVE_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "WAVE_ADJ_LIGHT",
        MSO_AUTO_SHAPE_TYPE.WAVE,
        {"adj1": 10_000, "adj2": 0},
    ),
    ShapeVariant(
        "WAVE_ADJ_SHIFT",
        MSO_AUTO_SHAPE_TYPE.WAVE,
        {"adj1": 12_500, "adj2": 40_000},
    ),
    ShapeVariant(
        "WAVE_ADJ_DEEP",
        MSO_AUTO_SHAPE_TYPE.WAVE,
        {"adj1": 30_000, "adj2": 0},
    ),
    ShapeVariant(
        "WAVE_ADJ_DEEP_SHIFT",
        MSO_AUTO_SHAPE_TYPE.WAVE,
        {"adj1": 30_000, "adj2": 40_000},
    ),
)

DOUBLE_WAVE_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "DOUBLE_WAVE_ADJ_LIGHT",
        MSO_AUTO_SHAPE_TYPE.DOUBLE_WAVE,
        {"adj1": 10_000, "adj2": 0},
    ),
    ShapeVariant(
        "DOUBLE_WAVE_ADJ_SHIFT",
        MSO_AUTO_SHAPE_TYPE.DOUBLE_WAVE,
        {"adj1": 12_500, "adj2": 40_000},
    ),
    ShapeVariant(
        "DOUBLE_WAVE_ADJ_DEEP",
        MSO_AUTO_SHAPE_TYPE.DOUBLE_WAVE,
        {"adj1": 30_000, "adj2": 0},
    ),
    ShapeVariant(
        "DOUBLE_WAVE_ADJ_DEEP_SHIFT",
        MSO_AUTO_SHAPE_TYPE.DOUBLE_WAVE,
        {"adj1": 30_000, "adj2": 40_000},
    ),
)

QUAD_ARROW_CALLOUT_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "QUAD_ARROW_CALLOUT_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW_CALLOUT,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000, "adj4": 15_000},
    ),
    ShapeVariant(
        "QUAD_ARROW_CALLOUT_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW_CALLOUT,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 35_000, "adj4": 35_000},
    ),
    ShapeVariant(
        "QUAD_ARROW_CALLOUT_ADJ_LONG",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW_CALLOUT,
        {"adj1": 20_000, "adj2": 50_000, "adj3": 25_000, "adj4": 50_000},
    ),
    ShapeVariant(
        "QUAD_ARROW_CALLOUT_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW_CALLOUT,
        {"adj1": 45_000, "adj2": 20_000, "adj3": 45_000, "adj4": 20_000},
    ),
)

CLOUD_CALLOUT_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "CLOUD_CALLOUT_ADJ_LEFT",
        MSO_AUTO_SHAPE_TYPE.CLOUD_CALLOUT,
        {"adj1": -20_000, "adj2": 30_000},
    ),
    ShapeVariant(
        "CLOUD_CALLOUT_ADJ_RIGHT",
        MSO_AUTO_SHAPE_TYPE.CLOUD_CALLOUT,
        {"adj1": 20_000, "adj2": 30_000},
    ),
    ShapeVariant(
        "CLOUD_CALLOUT_ADJ_LOW",
        MSO_AUTO_SHAPE_TYPE.CLOUD_CALLOUT,
        {"adj1": 0, "adj2": 80_000},
    ),
    ShapeVariant(
        "CLOUD_CALLOUT_ADJ_HIGH",
        MSO_AUTO_SHAPE_TYPE.CLOUD_CALLOUT,
        {"adj1": 0, "adj2": 10_000},
    ),
)

TEARDROP_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "TEARDROP_ADJ_LIGHT",
        MSO_AUTO_SHAPE_TYPE.TEAR,
        {"adj": 20_000},
    ),
    ShapeVariant(
        "TEARDROP_ADJ_DEFAULT",
        MSO_AUTO_SHAPE_TYPE.TEAR,
        {"adj": 50_000},
    ),
    ShapeVariant(
        "TEARDROP_ADJ_DEEP",
        MSO_AUTO_SHAPE_TYPE.TEAR,
        {"adj": 80_000},
    ),
    ShapeVariant(
        "TEARDROP_ADJ_SHARP",
        MSO_AUTO_SHAPE_TYPE.TEAR,
        {"adj": 100_000},
    ),
)

PIE_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "PIE_ADJ_SMALL",
        MSO_AUTO_SHAPE_TYPE.PIE,
        {"adj1": 3_000_000, "adj2": 12_000_000},
    ),
    ShapeVariant(
        "PIE_ADJ_HALF",
        MSO_AUTO_SHAPE_TYPE.PIE,
        {"adj1": 5_400_000, "adj2": 16_200_000},
    ),
    ShapeVariant(
        "PIE_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.PIE,
        {"adj1": 0, "adj2": 18_000_000},
    ),
    ShapeVariant(
        "PIE_ADJ_SLIVER",
        MSO_AUTO_SHAPE_TYPE.PIE,
        {"adj1": 9_000_000, "adj2": 11_000_000},
    ),
)

ARC_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "ARC_ADJ_SMALL",
        MSO_AUTO_SHAPE_TYPE.ARC,
        {"adj1": 3_000_000, "adj2": 12_000_000},
    ),
    ShapeVariant(
        "ARC_ADJ_HALF",
        MSO_AUTO_SHAPE_TYPE.ARC,
        {"adj1": 5_400_000, "adj2": 16_200_000},
    ),
    ShapeVariant(
        "ARC_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.ARC,
        {"adj1": 0, "adj2": 18_000_000},
    ),
    ShapeVariant(
        "ARC_ADJ_SLIVER",
        MSO_AUTO_SHAPE_TYPE.ARC,
        {"adj1": 9_000_000, "adj2": 11_000_000},
    ),
)

BLOCK_ARC_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "BLOCK_ARC_ADJ_NARROW",
        MSO_AUTO_SHAPE_TYPE.BLOCK_ARC,
        {"adj1": 12_000, "adj2": 8_500_000, "adj3": 17_000_000},
    ),
    ShapeVariant(
        "BLOCK_ARC_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.BLOCK_ARC,
        {"adj1": 35_000, "adj2": 3_000_000, "adj3": 13_000_000},
    ),
    ShapeVariant(
        "BLOCK_ARC_ADJ_RING",
        MSO_AUTO_SHAPE_TYPE.BLOCK_ARC,
        {"adj1": 50_000, "adj2": 0, "adj3": 21_600_000},
    ),
    ShapeVariant(
        "BLOCK_ARC_ADJ_OFFSET",
        MSO_AUTO_SHAPE_TYPE.BLOCK_ARC,
        {"adj1": 25_000, "adj2": 6_000_000, "adj3": 18_000_000},
    ),
)

LEFT_RIGHT_ARROW_CALLOUT_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "LEFT_RIGHT_ARROW_CALLOUT_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW_CALLOUT,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000, "adj4": 15_000},
    ),
    ShapeVariant(
        "LEFT_RIGHT_ARROW_CALLOUT_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW_CALLOUT,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 35_000, "adj4": 35_000},
    ),
    ShapeVariant(
        "LEFT_RIGHT_ARROW_CALLOUT_ADJ_LONG",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW_CALLOUT,
        {"adj1": 20_000, "adj2": 50_000, "adj3": 25_000, "adj4": 50_000},
    ),
    ShapeVariant(
        "LEFT_RIGHT_ARROW_CALLOUT_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW_CALLOUT,
        {"adj1": 45_000, "adj2": 20_000, "adj3": 45_000, "adj4": 20_000},
    ),
)

NOTCHED_RIGHT_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "NOTCHED_RIGHT_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW,
        {"adj1": 15_000, "adj2": 15_000},
    ),
    ShapeVariant(
        "NOTCHED_RIGHT_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW,
        {"adj1": 35_000, "adj2": 35_000},
    ),
    ShapeVariant(
        "NOTCHED_RIGHT_ARROW_ADJ_LONG",
        MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW,
        {"adj1": 20_000, "adj2": 50_000},
    ),
    ShapeVariant(
        "NOTCHED_RIGHT_ARROW_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW,
        {"adj1": 45_000, "adj2": 20_000},
    ),
)

BENT_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "BENT_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.BENT_ARROW,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000, "adj4": 35_000},
    ),
    ShapeVariant(
        "BENT_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.BENT_ARROW,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 35_000, "adj4": 50_000},
    ),
    ShapeVariant(
        "BENT_ARROW_ADJ_TALL",
        MSO_AUTO_SHAPE_TYPE.BENT_ARROW,
        {"adj1": 20_000, "adj2": 20_000, "adj3": 50_000, "adj4": 65_000},
    ),
    ShapeVariant(
        "BENT_ARROW_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.BENT_ARROW,
        {"adj1": 45_000, "adj2": 15_000, "adj3": 25_000, "adj4": 25_000},
    ),
)

QUAD_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "QUAD_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000},
    ),
    ShapeVariant(
        "QUAD_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 35_000},
    ),
    ShapeVariant(
        "QUAD_ARROW_ADJ_TALL",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW,
        {"adj1": 20_000, "adj2": 20_000, "adj3": 50_000},
    ),
    ShapeVariant(
        "QUAD_ARROW_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.QUAD_ARROW,
        {"adj1": 45_000, "adj2": 15_000, "adj3": 25_000},
    ),
)

LEFT_RIGHT_UP_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "LEFT_RIGHT_UP_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_UP_ARROW,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000},
    ),
    ShapeVariant(
        "LEFT_RIGHT_UP_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_UP_ARROW,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 35_000},
    ),
    ShapeVariant(
        "LEFT_RIGHT_UP_ARROW_ADJ_TALL",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_UP_ARROW,
        {"adj1": 20_000, "adj2": 20_000, "adj3": 50_000},
    ),
    ShapeVariant(
        "LEFT_RIGHT_UP_ARROW_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_UP_ARROW,
        {"adj1": 45_000, "adj2": 15_000, "adj3": 25_000},
    ),
)

UTURN_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "UTURN_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.U_TURN_ARROW,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000, "adj4": 35_000, "adj5": 65_000},
    ),
    ShapeVariant(
        "UTURN_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.U_TURN_ARROW,
        {"adj1": 35_000, "adj2": 30_000, "adj3": 35_000, "adj4": 50_000, "adj5": 85_000},
    ),
    ShapeVariant(
        "UTURN_ARROW_ADJ_SHALLOW",
        MSO_AUTO_SHAPE_TYPE.U_TURN_ARROW,
        {"adj1": 20_000, "adj2": 12_000, "adj3": 25_000, "adj4": 65_000, "adj5": 85_000},
    ),
    ShapeVariant(
        "UTURN_ARROW_ADJ_DEEP",
        MSO_AUTO_SHAPE_TYPE.U_TURN_ARROW,
        {"adj1": 25_000, "adj2": 45_000, "adj3": 10_000, "adj4": 25_000, "adj5": 70_000},
    ),
)

LEFT_UP_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "LEFT_UP_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.LEFT_UP_ARROW,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000},
    ),
    ShapeVariant(
        "LEFT_UP_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.LEFT_UP_ARROW,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 45_000},
    ),
    ShapeVariant(
        "LEFT_UP_ARROW_ADJ_LONG",
        MSO_AUTO_SHAPE_TYPE.LEFT_UP_ARROW,
        {"adj1": 20_000, "adj2": 20_000, "adj3": 50_000},
    ),
    ShapeVariant(
        "LEFT_UP_ARROW_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.LEFT_UP_ARROW,
        {"adj1": 40_000, "adj2": 15_000, "adj3": 25_000},
    ),
)

BENT_UP_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "BENT_UP_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.BENT_UP_ARROW,
        {"adj1": 15_000, "adj2": 15_000, "adj3": 15_000},
    ),
    ShapeVariant(
        "BENT_UP_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.BENT_UP_ARROW,
        {"adj1": 35_000, "adj2": 35_000, "adj3": 40_000},
    ),
    ShapeVariant(
        "BENT_UP_ARROW_ADJ_TALL",
        MSO_AUTO_SHAPE_TYPE.BENT_UP_ARROW,
        {"adj1": 25_000, "adj2": 15_000, "adj3": 50_000},
    ),
    ShapeVariant(
        "BENT_UP_ARROW_ADJ_DEEP",
        MSO_AUTO_SHAPE_TYPE.BENT_UP_ARROW,
        {"adj1": 20_000, "adj2": 45_000, "adj3": 20_000},
    ),
)

CIRCULAR_ARROW_VARIANTS: tuple[ShapeVariant, ...] = (
    ShapeVariant(
        "CIRCULAR_ARROW_ADJ_TIGHT",
        MSO_AUTO_SHAPE_TYPE.CIRCULAR_ARROW,
        {"adj1": -20_000, "adj5": 10_000},
    ),
    ShapeVariant(
        "CIRCULAR_ARROW_ADJ_WIDE",
        MSO_AUTO_SHAPE_TYPE.CIRCULAR_ARROW,
        {"adj1": 25_000, "adj5": 35_000},
    ),
    ShapeVariant(
        "CIRCULAR_ARROW_ADJ_SWEEP",
        MSO_AUTO_SHAPE_TYPE.CIRCULAR_ARROW,
        {"adj1": 45_000, "adj5": 15_000},
    ),
    ShapeVariant(
        "CIRCULAR_ARROW_ADJ_THICK",
        MSO_AUTO_SHAPE_TYPE.CIRCULAR_ARROW,
        {"adj1": 12_500, "adj5": 45_000},
    ),
)

SCENARIOS: dict[str, tuple[ShapeVariant, ...]] = {
    "curved-arrows": CURVED_ARROW_VARIANTS,
    "circular-arrow": CIRCULAR_ARROW_VARIANTS,
    "bent-up-arrow": BENT_UP_ARROW_VARIANTS,
    "left-up-arrow": LEFT_UP_ARROW_VARIANTS,
    "uturn-arrow": UTURN_ARROW_VARIANTS,
    "left-right-up-arrow": LEFT_RIGHT_UP_ARROW_VARIANTS,
    "quad-arrow": QUAD_ARROW_VARIANTS,
    "notched-right-arrow": NOTCHED_RIGHT_ARROW_VARIANTS,
    "bent-arrow": BENT_ARROW_VARIANTS,
    "up-down-arrow-callout": UP_DOWN_ARROW_CALLOUT_VARIANTS,
    "cloud-callout": CLOUD_CALLOUT_VARIANTS,
    "teardrop": TEARDROP_VARIANTS,
    "pie": PIE_VARIANTS,
    "arc": ARC_VARIANTS,
    "block-arc": BLOCK_ARC_VARIANTS,
    "left-right-arrow-callout": LEFT_RIGHT_ARROW_CALLOUT_VARIANTS,
    "wave": WAVE_VARIANTS,
    "double-wave": DOUBLE_WAVE_VARIANTS,
    "quad-arrow-callout": QUAD_ARROW_CALLOUT_VARIANTS,
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--scenario",
        choices=sorted(SCENARIOS),
        default="curved-arrows",
        help="Curated variant set to generate (default: curved-arrows)",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        required=True,
        help="Directory where the PPTX deck and manifest should be written",
    )
    parser.add_argument(
        "--deck-name",
        default=None,
        help="Optional output basename without extension; defaults to the scenario name",
    )
    return parser.parse_args()


def apply_adjustments(shape, adjustments: dict[str, int]) -> None:
    prst_geom = shape._element.spPr.prstGeom
    for child in list(prst_geom):
        if child.tag.endswith("avLst"):
            prst_geom.remove(child)

    av_lst = OxmlElement("a:avLst")
    prst_geom.insert(0, av_lst)
    for name, value in adjustments.items():
        gd = OxmlElement("a:gd")
        gd.set("name", name)
        gd.set("fmla", f"val {int(value)}")
        av_lst.append(gd)


def build_curated_slide(
    prs: Presentation, variants: Iterable[ShapeVariant]
) -> tuple[list[dict[str, object]], dict[str, float]]:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    margin_x = 0.4
    margin_y = 0.55
    cols = 4
    rows = 2
    cell_w = (10 - margin_x * 2) / cols
    cell_h = (7.5 - margin_y * 2) / rows
    size_w = 1.75
    size_h = 1.75
    fill_rgb = RGBColor(0x44, 0x72, 0xC4)
    line_rgb = RGBColor(0x20, 0x20, 0x20)

    entries: list[dict[str, object]] = []
    for idx, variant in enumerate(variants):
        col = idx % cols
        row = idx // cols
        cell_left = margin_x + col * cell_w
        cell_top = margin_y + row * cell_h
        left = cell_left + (cell_w - size_w) / 2
        top = cell_top + (cell_h - size_h) / 2

        shape = slide.shapes.add_shape(
            variant.shape_type,
            Inches(left),
            Inches(top),
            Inches(size_w),
            Inches(size_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(1.5)
        apply_adjustments(shape, variant.adjustments)

        entries.append(
            {
                "shape_name": variant.label,
                "base_shape_name": variant.shape_type.name,
                "slide_index": 0,
                "slot_index": idx,
                "crop_in": {
                    "left": round(cell_left, 4),
                    "top": round(cell_top, 4),
                    "width": round(cell_w, 4),
                    "height": round(cell_h, 4),
                },
                "shape_box_in": {
                    "left": round(left, 4),
                    "top": round(top, 4),
                    "width": size_w,
                    "height": size_h,
                },
                "adjustments": variant.adjustments,
            }
        )

    return entries, {
        "cols": cols,
        "rows": rows,
        "cell_w_in": cell_w,
        "cell_h_in": cell_h,
    }


def main() -> None:
    args = parse_args()
    variants = SCENARIOS[args.scenario]
    deck_name = args.deck_name or args.scenario.replace("-", "_")

    output_dir = args.output_dir.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / f"{deck_name}.pptx"
    manifest_path = output_dir / "manifest.json"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    entries, grid = build_curated_slide(prs, variants)
    prs.save(pptx_path)

    manifest = {
        "pptx": str(pptx_path),
        "scenario": args.scenario,
        "shape_count": len(entries),
        "slide_count": len(prs.slides),
        "grid": grid,
        "entries": entries,
    }
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(
        json.dumps(
            {
                "pptx": str(pptx_path),
                "manifest": str(manifest_path),
                "scenario": args.scenario,
                "shape_count": len(entries),
            },
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
