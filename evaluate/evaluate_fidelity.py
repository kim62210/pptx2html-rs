#!/usr/bin/env python3
"""
Fidelity evaluation for pptx2html-rs autoresearch pipeline.

This is the immutable scoring function -- the autoresearch equivalent of
prepare.py.  The LLM MUST NEVER modify this file; only humans may change
the scoring weights or metric definitions.

Composite score:
    fidelity = 0.40 * ssim + 0.25 * text_match + 0.25 * test_pass + 0.10 * perf
"""

from __future__ import annotations

import argparse
import json
import logging
import subprocess
import sys
import time
from pathlib import Path
from typing import NamedTuple

from bs4 import BeautifulSoup
from PIL import Image
from skimage.metrics import structural_similarity as ssim

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Score weights -- DO NOT change without human review
# ---------------------------------------------------------------------------
W_SSIM = 0.40
W_TEXT = 0.25
W_TEST = 0.25
W_PERF = 0.10

# Performance baseline: slides per second that maps to score 1.0
PERF_BASELINE_SPS = 50.0

# Default viewport matching standard slide dimensions (px)
VIEWPORT_WIDTH = 960
VIEWPORT_HEIGHT = 720


# ---------------------------------------------------------------------------
# Data containers
# ---------------------------------------------------------------------------
class SlideScore(NamedTuple):
    """Per-slide evaluation result."""

    pptx_name: str
    slide_idx: int
    ssim_score: float
    text_match: float


class EvalResult(NamedTuple):
    """Aggregate evaluation result."""

    ssim_score: float
    text_match_score: float
    test_pass_rate: float
    tests_passed: int
    tests_total: int
    perf_score: float
    slides_per_sec: float
    fidelity_score: float
    slide_scores: list[SlideScore]


# ---------------------------------------------------------------------------
# Metric: visual SSIM
# ---------------------------------------------------------------------------
def _compute_ssim(ref_path: Path, cand_path: Path) -> float:
    """Compute SSIM between reference and candidate PNG images.

    Both images are converted to grayscale and resized to match if needed.
    Returns a float in [0, 1].
    """
    ref_img = Image.open(ref_path).convert("L")
    cand_img = Image.open(cand_path).convert("L")

    # Resize candidate to match reference dimensions
    if ref_img.size != cand_img.size:
        cand_img = cand_img.resize(ref_img.size, Image.LANCZOS)

    import numpy as np

    ref_arr = np.array(ref_img)
    cand_arr = np.array(cand_img)

    score: float = ssim(ref_arr, cand_arr, data_range=255)
    return max(0.0, min(1.0, score))


# ---------------------------------------------------------------------------
# Metric: text content match
# ---------------------------------------------------------------------------
def _extract_text_from_html(html_path: Path) -> str:
    """Extract visible text content from an HTML file."""
    content = html_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(content, "lxml")

    # Remove script and style elements
    for tag in soup(["script", "style"]):
        tag.decompose()

    text = soup.get_text(separator=" ", strip=True)
    return text


def _extract_text_from_pptx(pptx_path: Path) -> str:
    """Extract text content from a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        texts.append(para_text.strip())
    return " ".join(texts)


def _text_similarity(reference: str, candidate: str) -> float:
    """Compute token-level Jaccard similarity between two texts.

    Returns a float in [0, 1].
    """
    if not reference and not candidate:
        return 1.0
    if not reference or not candidate:
        return 0.0

    ref_tokens = set(reference.lower().split())
    cand_tokens = set(candidate.lower().split())

    if not ref_tokens:
        return 1.0 if not cand_tokens else 0.0

    intersection = ref_tokens & cand_tokens
    union = ref_tokens | cand_tokens
    return len(intersection) / len(union) if union else 0.0


# ---------------------------------------------------------------------------
# Metric: cargo test pass rate
# ---------------------------------------------------------------------------
def _run_cargo_tests(project_root: Path) -> tuple[int, int]:
    """Run `cargo test --workspace` and return (passed, total).

    Parses the test summary line from cargo output.
    """
    try:
        result = subprocess.run(
            ["cargo", "test", "--workspace"],
            cwd=str(project_root),
            capture_output=True,
            text=True,
            timeout=300,
        )
    except FileNotFoundError as exc:
        logger.warning("cargo not found, skipping test metric: %s", exc)
        return 0, 0
    except subprocess.TimeoutExpired:
        logger.warning("cargo test timed out after 300s")
        return 0, 0

    output = result.stdout + result.stderr
    passed = 0
    total = 0

    for line in output.splitlines():
        # Pattern: "test result: ok. 145 passed; 0 failed; 0 ignored; ..."
        if line.strip().startswith("test result:"):
            parts = line.split()
            for i, part in enumerate(parts):
                if part == "passed;" and i > 0:
                    try:
                        passed += int(parts[i - 1])
                    except ValueError:
                        pass
                if part == "failed;" and i > 0:
                    try:
                        failed_count = int(parts[i - 1])
                        total += passed + failed_count
                    except ValueError:
                        pass

    # If parsing failed, check return code
    if total == 0 and passed == 0:
        if result.returncode == 0:
            logger.info("cargo test passed but could not parse counts")
            return 1, 1
        return 0, 1

    return passed, total


# ---------------------------------------------------------------------------
# Metric: conversion performance
# ---------------------------------------------------------------------------
def _measure_performance(
    project_root: Path, pptx_files: list[Path]
) -> float:
    """Measure conversion speed in slides per second.

    Converts each PPTX file and measures wall-clock time.
    Returns slides_per_second.
    """
    cli_bin = project_root / "target" / "release" / "pptx2html"
    if not cli_bin.exists():
        # Try building release binary
        logger.info("Building release binary for performance measurement...")
        try:
            subprocess.run(
                ["cargo", "build", "--release", "--package", "pptx2html-cli"],
                cwd=str(project_root),
                capture_output=True,
                timeout=300,
            )
        except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
            logger.warning("Failed to build release binary: %s", exc)
            return 0.0

    if not cli_bin.exists():
        logger.warning("Release binary not found at %s", cli_bin)
        return 0.0

    total_slides = 0
    total_time = 0.0

    for pptx_file in pptx_files:
        from pptx import Presentation

        try:
            prs = Presentation(str(pptx_file))
            slide_count = len(prs.slides)
        except Exception as exc:
            logger.warning("Could not read %s: %s", pptx_file.name, exc)
            continue

        start = time.monotonic()
        try:
            subprocess.run(
                [str(cli_bin), str(pptx_file), "-o", "/dev/null"],
                capture_output=True,
                timeout=60,
            )
        except subprocess.TimeoutExpired:
            logger.warning("Conversion timed out for %s", pptx_file.name)
            continue

        elapsed = time.monotonic() - start
        total_slides += slide_count
        total_time += elapsed

    if total_time == 0:
        return 0.0

    return total_slides / total_time


# ---------------------------------------------------------------------------
# PPTX -> HTML conversion
# ---------------------------------------------------------------------------
def _convert_pptx(
    project_root: Path, pptx_path: Path, output_dir: Path
) -> Path | None:
    """Convert a PPTX file to HTML using pptx2html CLI.

    Returns the path to the generated HTML file, or None on failure.
    """
    cli_bin = project_root / "target" / "release" / "pptx2html"
    if not cli_bin.exists():
        cli_bin = project_root / "target" / "debug" / "pptx2html"

    if not cli_bin.exists():
        logger.error("pptx2html binary not found. Run `cargo build` first.")
        return None

    output_html = output_dir / f"{pptx_path.stem}.html"

    try:
        result = subprocess.run(
            [str(cli_bin), str(pptx_path), "-o", str(output_html)],
            capture_output=True,
            text=True,
            timeout=60,
        )
    except subprocess.TimeoutExpired:
        logger.warning("Conversion timed out for %s", pptx_path.name)
        return None

    if result.returncode != 0:
        logger.warning(
            "Conversion failed for %s: %s",
            pptx_path.name,
            result.stderr[:200],
        )
        return None

    if not output_html.exists():
        logger.warning("Expected output not found: %s", output_html)
        return None

    return output_html


# ---------------------------------------------------------------------------
# HTML -> PNG screenshot
# ---------------------------------------------------------------------------
def _screenshot_html(html_path: Path, output_dir: Path) -> list[Path]:
    """Take screenshots of HTML slides using Playwright.

    Returns list of PNG paths, one per slide.
    """
    from candidate_render import render_html_to_pngs

    return render_html_to_pngs(
        html_path=html_path,
        output_dir=output_dir,
        viewport_width=VIEWPORT_WIDTH,
        viewport_height=VIEWPORT_HEIGHT,
    )


# ---------------------------------------------------------------------------
# Core evaluation
# ---------------------------------------------------------------------------
def evaluate(
    project_root: Path,
    golden_set_dir: Path | None = None,
    references_dir: Path | None = None,
    phase: str | None = None,
    verbose: bool = False,
) -> EvalResult:
    """Run the full fidelity evaluation.

    Args:
        project_root: Path to the pptx2html-rs project root.
        golden_set_dir: Directory containing golden PPTX files.
        references_dir: Directory containing LibreOffice reference PNGs.
        phase: If set, filter golden set to this category prefix.
        verbose: If True, log per-slide scores.

    Returns:
        EvalResult with all metric scores.
    """
    eval_dir = Path(__file__).resolve().parent

    if golden_set_dir is None:
        golden_set_dir = eval_dir / "golden_set"
    if references_dir is None:
        references_dir = eval_dir / "golden_references"

    # Collect PPTX files
    pptx_files = sorted(golden_set_dir.glob("*.pptx"))
    if phase:
        pptx_files = [f for f in pptx_files if f.stem.startswith(phase)]

    if not pptx_files:
        logger.error("No PPTX files found in %s", golden_set_dir)
        return EvalResult(
            ssim_score=0.0,
            text_match_score=0.0,
            test_pass_rate=0.0,
            tests_passed=0,
            tests_total=0,
            perf_score=0.0,
            slides_per_sec=0.0,
            fidelity_score=0.0,
            slide_scores=[],
        )

    logger.info("Found %d PPTX files for evaluation", len(pptx_files))

    # Temporary directories for conversion output
    import tempfile

    tmp_base = Path(tempfile.mkdtemp(prefix="eval_fidelity_"))
    html_dir = tmp_base / "html"
    candidate_dir = tmp_base / "candidates"
    html_dir.mkdir()
    candidate_dir.mkdir()

    # --- Visual SSIM & Text Match ---
    slide_scores: list[SlideScore] = []

    for pptx_file in pptx_files:
        ref_dir = references_dir / pptx_file.stem
        if not ref_dir.exists():
            logger.warning(
                "No reference renders for %s, skipping visual eval",
                pptx_file.name,
            )
            continue

        # Convert PPTX -> HTML
        html_path = _convert_pptx(project_root, pptx_file, html_dir)
        if html_path is None:
            continue

        # Screenshot HTML -> PNGs
        try:
            candidate_pngs = _screenshot_html(
                html_path, candidate_dir / pptx_file.stem
            )
        except Exception as exc:
            logger.warning(
                "Screenshot failed for %s: %s", pptx_file.name, exc
            )
            continue

        # Reference PNGs
        ref_pngs = sorted(ref_dir.glob("slide_*.png"))

        # SSIM per slide
        for idx, ref_png in enumerate(ref_pngs):
            if idx >= len(candidate_pngs):
                slide_scores.append(
                    SlideScore(
                        pptx_name=pptx_file.stem,
                        slide_idx=idx,
                        ssim_score=0.0,
                        text_match=0.0,
                    )
                )
                continue

            ssim_val = _compute_ssim(ref_png, candidate_pngs[idx])

            # Text match for this slide (aggregate at file level)
            text_match_val = 0.0
            if idx == 0:
                pptx_text = _extract_text_from_pptx(pptx_file)
                html_text = _extract_text_from_html(html_path)
                text_match_val = _text_similarity(pptx_text, html_text)

            slide_scores.append(
                SlideScore(
                    pptx_name=pptx_file.stem,
                    slide_idx=idx,
                    ssim_score=ssim_val,
                    text_match=text_match_val,
                )
            )

            if verbose:
                logger.info(
                    "  %s slide %d: SSIM=%.3f text=%.3f",
                    pptx_file.stem,
                    idx,
                    ssim_val,
                    text_match_val,
                )

    # Aggregate SSIM
    ssim_values = [s.ssim_score for s in slide_scores]
    avg_ssim = sum(ssim_values) / len(ssim_values) if ssim_values else 0.0

    # Aggregate text match (only first slide per PPTX carries the score)
    text_scores = [s.text_match for s in slide_scores if s.slide_idx == 0]
    avg_text = sum(text_scores) / len(text_scores) if text_scores else 0.0

    # --- cargo test ---
    logger.info("Running cargo test...")
    tests_passed, tests_total = _run_cargo_tests(project_root)
    test_rate = tests_passed / tests_total if tests_total > 0 else 0.0

    # --- Performance ---
    logger.info("Measuring conversion performance...")
    sps = _measure_performance(project_root, pptx_files)
    perf_score = min(1.0, sps / PERF_BASELINE_SPS)

    # --- Composite ---
    fidelity = (
        W_SSIM * avg_ssim
        + W_TEXT * avg_text
        + W_TEST * test_rate
        + W_PERF * perf_score
    )

    # Cleanup
    import shutil

    shutil.rmtree(tmp_base, ignore_errors=True)

    return EvalResult(
        ssim_score=avg_ssim,
        text_match_score=avg_text,
        test_pass_rate=test_rate,
        tests_passed=tests_passed,
        tests_total=tests_total,
        perf_score=perf_score,
        slides_per_sec=sps,
        fidelity_score=fidelity,
        slide_scores=slide_scores,
    )


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def _print_result(result: EvalResult) -> None:
    """Print evaluation result to stdout."""
    logger.info("=== Fidelity Evaluation ===")
    logger.info("SSIM Score:     %.3f", result.ssim_score)
    logger.info("Text Match:     %.3f", result.text_match_score)
    logger.info(
        "Test Pass Rate: %.3f (%d/%d)",
        result.test_pass_rate,
        result.tests_passed,
        result.tests_total,
    )
    logger.info(
        "Performance:    %.3f (%.1f slides/sec)",
        result.perf_score,
        result.slides_per_sec,
    )
    logger.info("---")
    logger.info("FIDELITY_SCORE: %.3f", result.fidelity_score)


def _write_json_result(result: EvalResult, output_path: Path) -> None:
    """Write evaluation result as JSON."""
    data = {
        "ssim_score": result.ssim_score,
        "text_match_score": result.text_match_score,
        "test_pass_rate": result.test_pass_rate,
        "tests_passed": result.tests_passed,
        "tests_total": result.tests_total,
        "perf_score": result.perf_score,
        "slides_per_sec": result.slides_per_sec,
        "fidelity_score": result.fidelity_score,
    }
    output_path.write_text(json.dumps(data, indent=2), encoding="utf-8")
    logger.info("JSON result written to %s", output_path)


def main() -> None:
    """CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Evaluate pptx2html-rs conversion fidelity"
    )
    parser.add_argument(
        "--project-root",
        type=Path,
        required=True,
        help="Path to pptx2html-rs project root",
    )
    parser.add_argument(
        "--golden-set",
        type=Path,
        default=None,
        help="Directory containing golden PPTX files",
    )
    parser.add_argument(
        "--references",
        type=Path,
        default=None,
        help="Directory containing reference PNG renders",
    )
    parser.add_argument(
        "--phase",
        type=str,
        default=None,
        help="Filter golden set by category prefix (e.g., 'color', 'text')",
    )
    parser.add_argument(
        "--output-json",
        type=Path,
        default=None,
        help="Write result as JSON to this path",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Show per-slide scores",
    )

    args = parser.parse_args()

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(message)s",
    )

    project_root = args.project_root.resolve()
    if not (project_root / "Cargo.toml").exists():
        logger.error("No Cargo.toml found at %s", project_root)
        sys.exit(1)

    result = evaluate(
        project_root=project_root,
        golden_set_dir=args.golden_set,
        references_dir=args.references,
        phase=args.phase,
        verbose=args.verbose,
    )

    _print_result(result)

    if args.output_json:
        _write_json_result(result, args.output_json)

    # Exit with non-zero if score is below threshold
    if result.fidelity_score < 0.01:
        sys.exit(2)


if __name__ == "__main__":
    main()
