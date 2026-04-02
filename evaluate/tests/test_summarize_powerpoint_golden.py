import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from evaluate.summarize_powerpoint_golden import summarize_powerpoint_golden_batch


class SummarizePowerPointGoldenBatchTests(unittest.TestCase):
    def test_reports_missing_metadata_and_missing_deck_output(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            self._create_pptx(golden_set / "sample_a.pptx", slide_count=2)
            self._create_pptx(golden_set / "sample_b.pptx", slide_count=1)

            deck_output = output_dir / "sample_a"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")

            summary = summarize_powerpoint_golden_batch(golden_set, output_dir)

            self.assertEqual(summary["golden_deck_count"], 2)
            self.assertEqual(summary["captured_deck_count"], 1)
            self.assertEqual(summary["missing_decks"], ["sample_b"])
            self.assertEqual(summary["missing_metadata"], ["sample_a"])
            self.assertEqual(summary["invalid_metadata"], [])
            self.assertEqual(summary["incomplete_slide_exports"], ["sample_a"])
            self.assertFalse(summary["manifest_present"])
            self.assertFalse(summary["manifest_deck_count_matches"])
            self.assertFalse(summary["manifest_slide_count_matches"])
            self.assertFalse(summary["evidence_ready_for_exact_promotion"])

    def test_reports_complete_batch_and_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            deck_path = golden_set / "sample.pptx"
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            (deck_output / "Slide2.PNG").write_bytes(b"png2")
            (deck_output / "metadata.json").write_text(
                json.dumps(
                    {
                        "powerpoint_version": "16.0.17726.20160",
                        "powerpoint_channel": "Current Channel",
                        "windows_version": "Windows 11 23H2",
                        "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "output_resolution": "960x540",
                        "golden_set_revision": "abc1234",
                        "capture_date": "2026-04-02",
                        "deck_name": "sample",
                        "slide_count": 2,
                    }
                ),
                encoding="utf-8",
            )
            (output_dir / "manifest.json").write_text(
                json.dumps(
                    {
                        "powerpoint_version": "16.0.17726.20160",
                        "powerpoint_channel": "Current Channel",
                        "windows_version": "Windows 11 23H2",
                        "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "output_resolution": "960x540",
                        "golden_set_revision": "abc1234",
                        "capture_date": "2026-04-02",
                        "deck_count": 1,
                        "total_slide_count": 2,
                        "decks": [
                            {"name": "sample", "slide_count": 2, "output_dir": "sample"}
                        ],
                    }
                ),
                encoding="utf-8",
            )

            summary = summarize_powerpoint_golden_batch(golden_set, output_dir)

            self.assertEqual(summary["golden_deck_count"], 1)
            self.assertEqual(summary["captured_deck_count"], 1)
            self.assertEqual(summary["missing_decks"], [])
            self.assertEqual(summary["missing_metadata"], [])
            self.assertEqual(summary["invalid_metadata"], [])
            self.assertEqual(summary["incomplete_slide_exports"], [])
            self.assertTrue(summary["manifest_present"])
            self.assertTrue(summary["manifest_deck_count_matches"])
            self.assertTrue(summary["manifest_slide_count_matches"])
            self.assertEqual(
                summary["batch_identity"]["golden_set_revision"], "abc1234"
            )
            self.assertTrue(summary["evidence_ready_for_exact_promotion"])

    def test_reports_manifest_count_mismatch(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            deck_path = golden_set / "sample.pptx"
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            (deck_output / "Slide2.PNG").write_bytes(b"png2")
            (deck_output / "metadata.json").write_text(
                json.dumps(
                    {
                        "powerpoint_version": "16.0.17726.20160",
                        "powerpoint_channel": "Current Channel",
                        "windows_version": "Windows 11 23H2",
                        "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "output_resolution": "960x540",
                        "golden_set_revision": "abc1234",
                        "capture_date": "2026-04-02",
                        "deck_name": "sample",
                        "slide_count": 2,
                    }
                ),
                encoding="utf-8",
            )
            (output_dir / "manifest.json").write_text(
                json.dumps(
                    {
                        "powerpoint_version": "16.0.17726.20160",
                        "powerpoint_channel": "Current Channel",
                        "windows_version": "Windows 11 23H2",
                        "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "output_resolution": "960x540",
                        "golden_set_revision": "abc1234",
                        "capture_date": "2026-04-02",
                        "deck_count": 99,
                        "total_slide_count": 99,
                        "decks": [
                            {"name": "sample", "slide_count": 2, "output_dir": "sample"}
                        ],
                    }
                ),
                encoding="utf-8",
            )

            summary = summarize_powerpoint_golden_batch(golden_set, output_dir)

            self.assertTrue(summary["manifest_present"])
            self.assertFalse(summary["manifest_deck_count_matches"])
            self.assertFalse(summary["manifest_slide_count_matches"])
            self.assertFalse(summary["evidence_ready_for_exact_promotion"])

    def test_reports_invalid_metadata_fields_and_bad_slide_names(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            deck_path = golden_set / "sample.pptx"
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            (deck_output / "Slide3.PNG").write_bytes(b"png3")
            (deck_output / "metadata.json").write_text(
                json.dumps(
                    {
                        "powerpoint_version": "",
                        "powerpoint_channel": "Current Channel",
                        "windows_version": "Windows 11 23H2",
                        "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "output_resolution": "960x540",
                        "golden_set_revision": "abc1234",
                        "capture_date": "2026-04-02",
                        "deck_name": "sample",
                        "slide_count": 2,
                    }
                ),
                encoding="utf-8",
            )

            summary = summarize_powerpoint_golden_batch(golden_set, output_dir)

            self.assertEqual(summary["invalid_metadata"], ["sample"])
            self.assertEqual(summary["incomplete_slide_exports"], ["sample"])
            self.assertFalse(summary["evidence_ready_for_exact_promotion"])

    def _create_pptx(self, path: Path, slide_count: int) -> None:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        while len(presentation.slides) < slide_count:
            presentation.slides.add_slide(blank_layout)
        presentation.save(path)


if __name__ == "__main__":
    unittest.main()
