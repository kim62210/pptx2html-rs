import contextlib
import io
import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from evaluate.powerpoint_evidence import main


TEXT_LAYOUT_DECKS = [
    "basic_text_08_narrow_box_autofit",
    "basic_text_09_mixed_font_paragraph",
    "basic_text_10_bodypr_fidelity",
    "basic_text_11_wrap_gate_sentence",
    "basic_text_12_wrap_gate_unbreakable",
    "basic_text_13_autofit_modes",
    "basic_text_14_complex_script_fonts",
    "basic_text_15_mixed_script_single_run",
    "basic_text_16_cjk_autofit_wrap_gate",
    "basic_text_17_indic_complex_script_fonts",
    "basic_text_18_emoji_cluster_segments",
]


class PowerPointEvidenceCliTests(unittest.TestCase):
    def test_ready_command_returns_nonzero_when_batch_is_incomplete(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(golden_set / "sample.pptx", slide_count=1)

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "ready",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                    ]
                )

            self.assertEqual(exit_code, 1)
            payload = json.loads(stdout.getvalue())
            self.assertFalse(payload["evidence_ready_for_exact_promotion"])

    def test_ready_command_returns_zero_when_batch_is_complete(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            deck_path = golden_set / "sample.pptx"
            self._create_pptx(deck_path, slide_count=1)
            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                scaffold_exit = main(
                    [
                        "scaffold",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                        "--powerpoint-version",
                        "16.0.17726.20160",
                        "--powerpoint-channel",
                        "Current Channel",
                        "--windows-version",
                        "Windows 11 23H2",
                        "--export-command",
                        "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "--output-resolution",
                        "960x540",
                        "--golden-set-revision",
                        "abc1234",
                        "--capture-date",
                        "2026-04-02",
                    ]
                )
            self.assertEqual(scaffold_exit, 0)

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "ready",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                    ]
                )

            self.assertEqual(exit_code, 0)
            payload = json.loads(stdout.getvalue())
            self.assertTrue(payload["evidence_ready_for_exact_promotion"])

    def test_summary_command_emits_json(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(golden_set / "sample.pptx", slide_count=1)

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "summary",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                    ]
                )

            self.assertEqual(exit_code, 0)
            payload = json.loads(stdout.getvalue())
            self.assertEqual(payload["golden_deck_count"], 1)
            self.assertEqual(payload["missing_decks"], ["sample"])

    def test_summary_command_writes_output_json_file(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            report_path = root / "summary.json"
            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(golden_set / "sample.pptx", slide_count=1)

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "summary",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                        "--output-json",
                        str(report_path),
                    ]
                )

            self.assertEqual(exit_code, 0)
            self.assertTrue(report_path.is_file())
            payload = json.loads(report_path.read_text(encoding="utf-8"))
            self.assertEqual(payload["golden_deck_count"], 1)

    def test_gate_command_returns_nonzero_when_required_family_decks_are_missing(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(
                golden_set / f"{TEXT_LAYOUT_DECKS[0]}.pptx", slide_count=1
            )

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "gate",
                        "--family",
                        "text-layout",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                    ]
                )

            self.assertEqual(exit_code, 1)
            payload = json.loads(stdout.getvalue())
            self.assertFalse(payload["family_ready_for_exact_promotion"])
            self.assertIn(TEXT_LAYOUT_DECKS[1], payload["missing_required_decks"])

    def test_gate_command_returns_zero_when_required_family_decks_are_ready(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            for deck_name in TEXT_LAYOUT_DECKS:
                deck_path = golden_set / f"{deck_name}.pptx"
                self._create_pptx(deck_path, slide_count=1)
                deck_output = output_dir / deck_name
                deck_output.mkdir()
                (deck_output / "Slide1.PNG").write_bytes(b"png1")

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                scaffold_exit = main(
                    [
                        "scaffold",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                        "--powerpoint-version",
                        "16.0.17726.20160",
                        "--powerpoint-channel",
                        "Current Channel",
                        "--windows-version",
                        "Windows 11 23H2",
                        "--export-command",
                        "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "--output-resolution",
                        "960x540",
                        "--golden-set-revision",
                        "abc1234",
                        "--capture-date",
                        "2026-04-02",
                    ]
                )
            self.assertEqual(scaffold_exit, 0)

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "gate",
                        "--family",
                        "text-layout",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                    ]
                )

            self.assertEqual(exit_code, 0)
            payload = json.loads(stdout.getvalue())
            self.assertTrue(payload["family_ready_for_exact_promotion"])
            self.assertEqual(payload["missing_required_decks"], [])

    def test_gate_command_ignores_invalid_metadata_outside_family_bundle(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            golden_set.mkdir()
            output_dir.mkdir()

            for deck_name in TEXT_LAYOUT_DECKS:
                deck_path = golden_set / f"{deck_name}.pptx"
                self._create_pptx(deck_path, slide_count=1)
                deck_output = output_dir / deck_name
                deck_output.mkdir()
                (deck_output / "Slide1.PNG").write_bytes(b"png1")

            extra_deck = golden_set / "shapes_01_rectangles.pptx"
            self._create_pptx(extra_deck, slide_count=1)
            extra_output = output_dir / "shapes_01_rectangles"
            extra_output.mkdir()
            (extra_output / "Slide1.PNG").write_bytes(b"png1")

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                scaffold_exit = main(
                    [
                        "scaffold",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                        "--powerpoint-version",
                        "16.0.17726.20160",
                        "--powerpoint-channel",
                        "Current Channel",
                        "--windows-version",
                        "Windows 11 23H2",
                        "--export-command",
                        "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "--output-resolution",
                        "960x540",
                        "--golden-set-revision",
                        "abc1234",
                        "--capture-date",
                        "2026-04-02",
                    ]
                )
            self.assertEqual(scaffold_exit, 0)

            (extra_output / "metadata.json").write_text(
                json.dumps(
                    {
                        "powerpoint_version": "",
                        "powerpoint_channel": "Current Channel",
                        "windows_version": "Windows 11 23H2",
                        "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                        "output_resolution": "960x540",
                        "golden_set_revision": "abc1234",
                        "capture_date": "2026-04-02",
                        "deck_name": "shapes_01_rectangles",
                        "slide_count": 1,
                    }
                ),
                encoding="utf-8",
            )

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "gate",
                        "--family",
                        "text-layout",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                    ]
                )

            self.assertEqual(exit_code, 0)
            payload = json.loads(stdout.getvalue())
            self.assertTrue(payload["family_ready_for_exact_promotion"])
            self.assertEqual(payload["missing_required_decks"], [])

    def test_gate_command_writes_output_json_file_even_on_failure(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            report_path = root / "gate.json"
            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(
                golden_set / f"{TEXT_LAYOUT_DECKS[0]}.pptx", slide_count=1
            )

            stdout = io.StringIO()
            with contextlib.redirect_stdout(stdout):
                exit_code = main(
                    [
                        "gate",
                        "--family",
                        "text-layout",
                        "--golden-set-dir",
                        str(golden_set),
                        "--output-dir",
                        str(output_dir),
                        "--output-json",
                        str(report_path),
                    ]
                )

            self.assertEqual(exit_code, 1)
            self.assertTrue(report_path.is_file())
            payload = json.loads(report_path.read_text(encoding="utf-8"))
            self.assertFalse(payload["family_ready_for_exact_promotion"])

    def _create_pptx(self, path: Path, slide_count: int) -> None:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        while len(presentation.slides) < slide_count:
            presentation.slides.add_slide(blank_layout)
        presentation.save(path)


if __name__ == "__main__":
    unittest.main()
