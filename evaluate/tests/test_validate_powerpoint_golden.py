import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from evaluate.validate_powerpoint_golden import (
    ValidationError,
    validate_powerpoint_golden_batch,
)


class ValidatePowerPointGoldenBatchTests(unittest.TestCase):
    def test_rejects_missing_metadata_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            deck_path = golden_set / "sample.pptx"

            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            (deck_output / "Slide2.PNG").write_bytes(b"png2")

            with self.assertRaisesRegex(ValidationError, "metadata.json"):
                validate_powerpoint_golden_batch(golden_set, output_dir)

    def test_rejects_missing_slide_export(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            deck_path = golden_set / "sample.pptx"

            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            self._write_metadata(deck_output / "metadata.json")

            with self.assertRaisesRegex(ValidationError, "Slide2.PNG"):
                validate_powerpoint_golden_batch(golden_set, output_dir)

    def test_accepts_complete_batch_and_returns_summary(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            deck_path = golden_set / "sample.pptx"

            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            (deck_output / "Slide2.PNG").write_bytes(b"png2")
            self._write_metadata(deck_output / "metadata.json")

            summary = validate_powerpoint_golden_batch(golden_set, output_dir)

            self.assertEqual(summary["deck_count"], 1)
            self.assertEqual(summary["slide_image_count"], 2)
            self.assertEqual(summary["validated_decks"], ["sample"])

    def _create_pptx(self, path: Path, slide_count: int) -> None:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        while len(presentation.slides) < slide_count:
            presentation.slides.add_slide(blank_layout)
        presentation.save(path)

    def _write_metadata(self, path: Path) -> None:
        path.write_text(
            json.dumps(
                {
                    "powerpoint_version": "16.0.17726.20160",
                    "powerpoint_channel": "Current Channel",
                    "windows_version": "Windows 11 23H2",
                    "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
                    "output_resolution": "960x540",
                    "golden_set_revision": "abc1234",
                    "capture_date": "2026-04-02",
                }
            ),
            encoding="utf-8",
        )


if __name__ == "__main__":
    unittest.main()
