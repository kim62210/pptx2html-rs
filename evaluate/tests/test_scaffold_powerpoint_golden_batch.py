import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from evaluate.scaffold_powerpoint_golden_batch import (
    ScaffoldError,
    scaffold_powerpoint_golden_batch,
)


class ScaffoldPowerPointGoldenBatchTests(unittest.TestCase):
    def test_rejects_missing_export_directory_for_deck(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            deck_path = golden_set / "sample.pptx"

            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(deck_path, slide_count=2)

            with self.assertRaisesRegex(ScaffoldError, "sample"):
                scaffold_powerpoint_golden_batch(
                    golden_set,
                    output_dir,
                    metadata=self._metadata(),
                )

    def test_writes_metadata_files_and_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            root = Path(tmpdir)
            golden_set = root / "golden_set"
            output_dir = root / "powerpoint_golden"
            deck_path = golden_set / "sample.pptx"

            golden_set.mkdir()
            output_dir.mkdir()
            self._create_pptx(deck_path, slide_count=2)

            deck_output = output_dir / "sample"
            deck_output.mkdir()
            (deck_output / "Slide1.PNG").write_bytes(b"png1")
            (deck_output / "Slide2.PNG").write_bytes(b"png2")

            summary = scaffold_powerpoint_golden_batch(
                golden_set,
                output_dir,
                metadata=self._metadata(),
            )

            metadata_path = deck_output / "metadata.json"
            manifest_path = output_dir / "manifest.json"
            self.assertTrue(metadata_path.is_file())
            self.assertTrue(manifest_path.is_file())

            deck_metadata = json.loads(metadata_path.read_text(encoding="utf-8"))
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

            self.assertEqual(deck_metadata["deck_name"], "sample")
            self.assertEqual(deck_metadata["slide_count"], 2)
            self.assertEqual(manifest["deck_count"], 1)
            self.assertEqual(manifest["total_slide_count"], 2)
            self.assertEqual(manifest["decks"][0]["name"], "sample")
            self.assertEqual(manifest["decks"][0]["slide_count"], 2)
            self.assertEqual(summary["deck_count"], 1)
            self.assertEqual(summary["slide_image_count"], 2)

    def _create_pptx(self, path: Path, slide_count: int) -> None:
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        while len(presentation.slides) < slide_count:
            presentation.slides.add_slide(blank_layout)
        presentation.save(path)

    def _metadata(self) -> dict[str, str]:
        return {
            "powerpoint_version": "16.0.17726.20160",
            "powerpoint_channel": "Current Channel",
            "windows_version": "Windows 11 23H2",
            "export_command": "pwsh -File ./reference_render_powerpoint.ps1 -InputDir ./golden_set -OutputDir ./powerpoint_golden",
            "output_resolution": "960x540",
            "golden_set_revision": "abc1234",
            "capture_date": "2026-04-02",
        }


if __name__ == "__main__":
    unittest.main()
