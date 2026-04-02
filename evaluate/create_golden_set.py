#!/usr/bin/env python3
"""
Generate golden PPTX test files for the fidelity evaluation pipeline.

Creates a golden PPTX corpus across 10 categories using python-pptx.
Each file exercises specific PPTX features that pptx2html-rs must handle.

    Categories:
        basic_text  - Text formatting: bold, italic, size, color, underline, bodyPr fidelity
    shapes      - Preset shapes: rectangle, ellipse, arrows, stars
    theme_colors - Theme color palette with tint/shade/lumMod modifiers
    tables      - Cell merge, borders, background colors, alignment
    images      - Embedded raster images with various placements
    gradients   - Gradient fills: linear, radial, multi-stop
    groups      - Grouped shapes, rotation, nested groups
    layouts     - Master/layout inheritance, backgrounds, placeholders
    bullets     - Bullet lists, numbered lists, indentation levels
    mixed       - Complex slides combining multiple element types
"""

from __future__ import annotations

import argparse
import io
import logging
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)

# Slide dimensions (standard 10x7.5 inches)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------
def _new_presentation() -> Presentation:
    """Create a blank presentation with standard dimensions."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _add_blank_slide(prs: Presentation) -> object:
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def _create_test_image(
    width: int = 200, height: int = 150, color: str = "red"
) -> bytes:
    """Create a simple test PNG image in memory."""
    img = PILImage.new("RGB", (width, height), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Category: basic_text
# ---------------------------------------------------------------------------
def _create_basic_text(output_dir: Path) -> list[Path]:
    """Generate the full basic_text fixture family for text fidelity checks."""
    files: list[Path] = []

    # 1. Bold/Italic/Underline combinations
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    styles = [
        ("Bold text", True, False, False),
        ("Italic text", False, True, False),
        ("Underlined text", False, False, True),
        ("Bold italic", True, True, False),
        ("Bold italic underline", True, True, True),
    ]
    for text, bold, italic, underline in styles:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        run.font.size = Pt(24)

    path = output_dir / "basic_text_01_styles.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Font sizes
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True

    for size in [10, 14, 18, 24, 32, 48, 72]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"Size {size}pt"
        run.font.size = Pt(size)

    path = output_dir / "basic_text_02_sizes.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Font colors
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    colors = [
        ("Red", RGBColor(0xFF, 0x00, 0x00)),
        ("Green", RGBColor(0x00, 0x80, 0x00)),
        ("Blue", RGBColor(0x00, 0x00, 0xFF)),
        ("Orange", RGBColor(0xFF, 0xA5, 0x00)),
        ("Purple", RGBColor(0x80, 0x00, 0x80)),
    ]
    for text, color in colors:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{text} colored text"
        run.font.color.rgb = color
        run.font.size = Pt(28)

    path = output_dir / "basic_text_03_colors.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Paragraph alignment
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True

    alignments = [
        ("Left aligned text", PP_ALIGN.LEFT),
        ("Center aligned text", PP_ALIGN.CENTER),
        ("Right aligned text", PP_ALIGN.RIGHT),
        (
            "Justified text with enough words to show justification effect clearly",
            PP_ALIGN.JUSTIFY,
        ),
    ]
    for text, align in alignments:
        p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(20)

    path = output_dir / "basic_text_04_alignment.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Multi-paragraph with spacing
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(6))
    tf = txBox.text_frame
    tf.word_wrap = True

    paragraphs = [
        "First paragraph with normal spacing.",
        "Second paragraph after space.",
        "Third paragraph with different style.",
        "Fourth paragraph continues the content.",
        "Fifth and final paragraph of text.",
    ]
    for i, text in enumerate(paragraphs):
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(18)
        if i % 2 == 0:
            run.font.bold = True

    path = output_dir / "basic_text_05_paragraphs.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Theme fallback candidate"
    run.font.size = Pt(24)
    run.font.name = "Aptos"
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Explicit alternate font"
    run2.font.size = Pt(24)
    run2.font.name = "Calibri"

    path = output_dir / "basic_text_06_font_fallback.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(1.2), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Vertical text sample"
    run.font.size = Pt(20)
    body_pr = txBox.text_frame._txBody.bodyPr
    body_pr.set("vert", "vert")

    path = output_dir / "basic_text_07_vertical_text.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(1.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "A very narrow text box with a long sentence to stress line breaking, "
        "wrapping, and auto-fit behavior in downstream renderers."
    )
    run.font.size = Pt(18)

    path = output_dir / "basic_text_08_narrow_box_autofit.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Latin "
    run.font.name = "Calibri"
    run.font.size = Pt(24)
    run2 = p.add_run()
    run2.text = "한글 "
    run2.font.name = "Malgun Gothic"
    run2.font.size = Pt(24)
    run3 = p.add_run()
    run3.text = "Mixed "
    run3.font.name = "Arial"
    run3.font.size = Pt(24)
    run4 = p.add_run()
    run4.text = "日本語"
    run4.font.name = "Yu Gothic"
    run4.font.size = Pt(24)

    path = output_dir / "basic_text_09_mixed_font_paragraph.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(2.0), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    body_pr = txBox.text_frame._txBody.bodyPr
    body_pr.set("anchorCtr", "1")
    body_pr.set("rot", "5400000")
    body_pr.set("vert", "vert270")
    body_pr.set("lIns", "0")
    body_pr.set("rIns", "0")
    body_pr.set("tIns", "0")
    body_pr.set("bIns", "0")

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "BodyPr fidelity"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.underline = True
    run.font.color.rgb = RGBColor(0x41, 0x72, 0xC4)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "NoWrapVerticalCenterRotate"
    run2.font.size = Pt(16)

    path = output_dir / "basic_text_10_bodypr_fidelity.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(2.4), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "This sentence should wrap at spaces before emergency breaks are needed."
    run.font.size = Pt(18)

    path = output_dir / "basic_text_11_wrap_gate_sentence.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(1.5), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "SupercalifragilisticexpialidociousWithoutSpaces"
    run.font.size = Pt(18)

    path = output_dir / "basic_text_12_wrap_gate_unbreakable.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    labels = [
        ("No autofit", Inches(0.7), None),
        ("Normal autofit shrink", Inches(3.4), "norm"),
        ("Shape autofit grows", Inches(6.1), "shape"),
    ]
    for label, left, mode in labels:
        txBox = slide.shapes.add_textbox(left, Inches(1.0), Inches(2.0), Inches(3.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        body_pr = txBox.text_frame._txBody.bodyPr
        if mode == "norm":
            body_pr.append(
                parse_xml(
                    r'<a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" fontScale="65000" lnSpcReduction="15000"/>'
                )
            )
        elif mode == "shape":
            body_pr.append(
                parse_xml(
                    r'<a:spAutoFit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                )
            )

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}: a narrow text box used to compare wrapping and shrink behavior across autofit modes."
        run.font.size = Pt(20)

    path = output_dir / "basic_text_13_autofit_modes.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.2), Inches(2.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "مرحبا بالعالم"
    run.font.name = "Amiri"
    run.font.size = Pt(24)
    run2 = p.add_run()
    run2.text = " / שלום עולם"
    run2.font.name = "Noto Sans Hebrew"
    run2.font.size = Pt(24)

    p2 = tf.add_paragraph()
    run3 = p2.add_run()
    run3.text = "Arabic and Hebrew runs should keep their complex-script typefaces."
    run3.font.size = Pt(18)

    path = output_dir / "basic_text_14_complex_script_fonts.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.2), Inches(2.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Hello مرحبا world שלום"
    run.font.name = "Calibri"
    run.font.size = Pt(24)
    body_pr = txBox.text_frame._txBody.bodyPr
    body_pr.set("rtlCol", "0")

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "One run mixes Latin, Arabic, and Hebrew to validate span segmentation."
    run2.font.size = Pt(18)

    path = output_dir / "basic_text_15_mixed_script_single_run.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(2.0), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    body_pr = txBox.text_frame._txBody.bodyPr
    body_pr.append(
        parse_xml(
            r'<a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" fontScale="70000"/>'
        )
    )
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "자동줄바꿈이가능한한글문장은긴토큰처럼취급되면안됩니다"
    run.font.size = Pt(18)

    path = output_dir / "basic_text_16_cjk_autofit_wrap_gate.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.2), Inches(2.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "नमस्ते दुनिया"
    run.font.name = "Nirmala UI"
    run.font.size = Pt(24)
    run2 = p.add_run()
    run2.text = " / สวัสดีชาวโลก"
    run2.font.name = "Leelawadee UI"
    run2.font.size = Pt(24)

    p2 = tf.add_paragraph()
    run3 = p2.add_run()
    run3.text = "Indic and Thai runs should keep complex-script typefaces instead of latin fallbacks."
    run3.font.size = Pt(18)

    path = output_dir / "basic_text_17_indic_complex_script_fonts.pptx"
    prs.save(str(path))
    files.append(path)

    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.2), Inches(2.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "A👩‍💻B"
    run.font.name = "Segoe UI Emoji"
    run.font.size = Pt(24)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = (
        "Emoji ZWJ clusters should stay intact and receive emoji-aware font selection."
    )
    run2.font.size = Pt(18)

    path = output_dir / "basic_text_18_emoji_cluster_segments.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: shapes
# ---------------------------------------------------------------------------
def _create_shapes(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing preset shapes."""
    files: list[Path] = []

    # 1. Basic rectangles with fills
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    colors_list = [
        RGBColor(0xFF, 0x00, 0x00),
        RGBColor(0x00, 0xFF, 0x00),
        RGBColor(0x00, 0x00, 0xFF),
    ]
    for i, color in enumerate(colors_list):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1 + i * 3),
            Inches(2),
            Inches(2),
            Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    path = output_dir / "shapes_01_rectangles.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Ellipses and circles
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    for i in range(4):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5 + i * 2.5),
            Inches(2),
            Inches(2),
            Inches(2 - i * 0.3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(50 + i * 50, 100, 200 - i * 40)

    path = output_dir / "shapes_02_ellipses.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Arrows and lines
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    arrow_shapes = [
        MSO_SHAPE.RIGHT_ARROW,
        MSO_SHAPE.LEFT_ARROW,
        MSO_SHAPE.UP_ARROW,
        MSO_SHAPE.DOWN_ARROW,
    ]
    positions = [
        (Inches(1), Inches(1)),
        (Inches(5), Inches(1)),
        (Inches(1), Inches(4)),
        (Inches(5), Inches(4)),
    ]
    for shape_type, (left, top) in zip(arrow_shapes, positions):
        shape = slide.shapes.add_shape(shape_type, left, top, Inches(3), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x33, 0x66, 0x99)

    path = output_dir / "shapes_03_arrows.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Stars and polygons
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    star_shapes = [
        MSO_SHAPE.STAR_4_POINT,
        MSO_SHAPE.STAR_5_POINT,
        MSO_SHAPE.STAR_6_POINT,
        MSO_SHAPE.STAR_8_POINT,
    ]
    for i, star in enumerate(star_shapes):
        shape = slide.shapes.add_shape(
            star,
            Inches(0.5 + i * 2.5),
            Inches(2),
            Inches(2),
            Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0x00)

    path = output_dir / "shapes_04_stars.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Shapes with text inside
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    shape_specs = [
        (MSO_SHAPE.ROUNDED_RECTANGLE, "Rounded Rect"),
        (MSO_SHAPE.DIAMOND, "Diamond"),
        (MSO_SHAPE.HEXAGON, "Hexagon"),
        (MSO_SHAPE.PENTAGON, "Pentagon"),
    ]
    for i, (shape_type, label) in enumerate(shape_specs):
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(0.5 + i * 2.5),
            Inches(2),
            Inches(2),
            Inches(2.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        shape.text_frame.paragraphs[0].text = label
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in shape.text_frame.paragraphs[0].runs:
            run.font.size = Pt(14)

    path = output_dir / "shapes_05_with_text.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: theme_colors
# ---------------------------------------------------------------------------
def _create_theme_colors(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing theme color usage."""
    files: list[Path] = []

    # Theme color indices: dk1, lt1, dk2, lt2, accent1-6, hlink, folHlink
    # We test with explicit RGB since python-pptx theme color support is limited

    # 1. 12 theme-representative colors as shape fills
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    theme_colors = [
        ("dk1", RGBColor(0x00, 0x00, 0x00)),
        ("lt1", RGBColor(0xFF, 0xFF, 0xFF)),
        ("dk2", RGBColor(0x44, 0x54, 0x6A)),
        ("lt2", RGBColor(0xE7, 0xE6, 0xE6)),
        ("accent1", RGBColor(0x41, 0x72, 0xC4)),
        ("accent2", RGBColor(0xED, 0x7D, 0x31)),
        ("accent3", RGBColor(0xA5, 0xA5, 0xA5)),
        ("accent4", RGBColor(0xFB, 0xC0, 0x2D)),
        ("accent5", RGBColor(0x5B, 0x9B, 0xD5)),
        ("accent6", RGBColor(0x70, 0xAD, 0x47)),
        ("hlink", RGBColor(0x05, 0x63, 0xC1)),
        ("folHlink", RGBColor(0x95, 0x4F, 0x72)),
    ]
    for i, (name, color) in enumerate(theme_colors):
        row = i // 6
        col = i % 6
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + col * 1.5),
            Inches(1 + row * 3),
            Inches(1.2),
            Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        tf = shape.text_frame
        tf.paragraphs[0].text = name
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(10)
            # White text on dark backgrounds
            if name in ("dk1", "dk2", "accent1", "hlink", "folHlink"):
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    path = output_dir / "theme_colors_01_palette.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Tint variations of accent1
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    base_r, base_g, base_b = 0x41, 0x72, 0xC4
    tint_levels = [0.0, 0.2, 0.4, 0.6, 0.8]
    for i, tint in enumerate(tint_levels):
        r = int(base_r + (255 - base_r) * tint)
        g = int(base_g + (255 - base_g) * tint)
        b = int(base_b + (255 - base_b) * tint)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + i * 1.8),
            Inches(2),
            Inches(1.5),
            Inches(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        tf = shape.text_frame
        tf.paragraphs[0].text = f"tint {tint:.0%}"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    path = output_dir / "theme_colors_02_tint.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Shade variations of accent1
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    shade_levels = [0.0, 0.2, 0.4, 0.6, 0.8]
    for i, shade in enumerate(shade_levels):
        r = int(base_r * (1 - shade))
        g = int(base_g * (1 - shade))
        b = int(base_b * (1 - shade))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + i * 1.8),
            Inches(2),
            Inches(1.5),
            Inches(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        tf = shape.text_frame
        tf.paragraphs[0].text = f"shade {shade:.0%}"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in tf.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    path = output_dir / "theme_colors_03_shade.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. All 6 accent colors as text
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    tf = txBox.text_frame
    accent_colors = [
        ("Accent 1", RGBColor(0x41, 0x72, 0xC4)),
        ("Accent 2", RGBColor(0xED, 0x7D, 0x31)),
        ("Accent 3", RGBColor(0xA5, 0xA5, 0xA5)),
        ("Accent 4", RGBColor(0xFB, 0xC0, 0x2D)),
        ("Accent 5", RGBColor(0x5B, 0x9B, 0xD5)),
        ("Accent 6", RGBColor(0x70, 0xAD, 0x47)),
    ]
    for text, color in accent_colors:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"{text}: The quick brown fox"
        run.font.color.rgb = color
        run.font.size = Pt(28)
        run.font.bold = True

    path = output_dir / "theme_colors_04_accent_text.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Mixed theme colors with background
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    # Dark background shape
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0),
        Emu(0),
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    bg_shape.line.fill.background()

    # Light text boxes on dark background
    light_colors = [
        RGBColor(0xFF, 0xFF, 0xFF),
        RGBColor(0xE7, 0xE6, 0xE6),
        RGBColor(0x5B, 0x9B, 0xD5),
        RGBColor(0x70, 0xAD, 0x47),
    ]
    for i, color in enumerate(light_colors):
        txBox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1 + i * 1.5),
            Inches(8),
            Inches(1),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"Light text on dark background (variant {i + 1})"
        run.font.color.rgb = color
        run.font.size = Pt(24)

    path = output_dir / "theme_colors_05_dark_bg.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: tables
# ---------------------------------------------------------------------------
def _create_tables(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing table features."""
    files: list[Path] = []

    # 1. Simple table with header
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    rows, cols = 5, 4
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(8), Inches(5)
    )
    table = table_shape.table

    # Header row
    headers = ["Name", "Age", "City", "Score"]
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)

    # Data rows
    data = [
        ["Alice", "30", "New York", "95"],
        ["Bob", "25", "London", "87"],
        ["Charlie", "35", "Tokyo", "92"],
        ["Diana", "28", "Berlin", "88"],
    ]
    for i, row_data in enumerate(data):
        for j, value in enumerate(row_data):
            table.cell(i + 1, j).text = value

    path = output_dir / "tables_01_simple.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Table with colored cells
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    table_shape = slide.shapes.add_table(
        4, 3, Inches(1), Inches(1.5), Inches(8), Inches(4)
    )
    table = table_shape.table

    cell_colors = [
        [
            RGBColor(0x41, 0x72, 0xC4),
            RGBColor(0x5B, 0x9B, 0xD5),
            RGBColor(0x9D, 0xC3, 0xE6),
        ],
        [
            RGBColor(0xED, 0x7D, 0x31),
            RGBColor(0xF4, 0xB1, 0x83),
            RGBColor(0xFA, 0xD8, 0xC1),
        ],
        [
            RGBColor(0x70, 0xAD, 0x47),
            RGBColor(0xA9, 0xD1, 0x8E),
            RGBColor(0xD4, 0xE8, 0xC7),
        ],
        [
            RGBColor(0xFB, 0xC0, 0x2D),
            RGBColor(0xFD, 0xDB, 0x7D),
            RGBColor(0xFE, 0xED, 0xBE),
        ],
    ]
    for i in range(4):
        for j in range(3):
            cell = table.cell(i, j)
            cell.text = f"R{i}C{j}"
            cell.fill.solid()
            cell.fill.fore_color.rgb = cell_colors[i][j]

    path = output_dir / "tables_02_colored.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Table with merged cells
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    table_shape = slide.shapes.add_table(
        4, 4, Inches(1), Inches(1.5), Inches(8), Inches(4)
    )
    table = table_shape.table

    # Merge top row across all columns
    table.cell(0, 0).merge(table.cell(0, 3))
    table.cell(0, 0).text = "Merged Header Row"
    for paragraph in table.cell(0, 0).text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.bold = True

    # Merge first column rows 1-3
    table.cell(1, 0).merge(table.cell(3, 0))
    table.cell(1, 0).text = "Merged\nColumn"
    table.cell(1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Fill remaining cells
    for i in range(1, 4):
        for j in range(1, 4):
            table.cell(i, j).text = f"({i},{j})"

    path = output_dir / "tables_03_merged.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Table with different text alignment
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    table_shape = slide.shapes.add_table(
        3, 3, Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    )
    table = table_shape.table

    alignments = [
        [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT],
        [PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.LEFT],
        [PP_ALIGN.RIGHT, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
    ]
    v_anchors = [MSO_ANCHOR.TOP, MSO_ANCHOR.MIDDLE, MSO_ANCHOR.BOTTOM]

    for i in range(3):
        for j in range(3):
            cell = table.cell(i, j)
            cell.text = f"Align ({i},{j})"
            cell.text_frame.paragraphs[0].alignment = alignments[i][j]
            cell.vertical_anchor = v_anchors[i]

    path = output_dir / "tables_04_alignment.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Large table with many rows
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    rows, cols = 10, 5
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(0.5), Inches(9), Inches(6.5)
    )
    table = table_shape.table

    headers = ["ID", "Product", "Qty", "Price", "Total"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2E, 0x75, 0xB6)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i in range(1, rows):
        table.cell(i, 0).text = str(i)
        table.cell(i, 1).text = f"Product {chr(64 + i)}"
        table.cell(i, 2).text = str(i * 10)
        table.cell(i, 3).text = f"${i * 9.99:.2f}"
        table.cell(i, 4).text = f"${i * 10 * i * 9.99:.2f}"

        if i % 2 == 0:
            for j in range(cols):
                table.cell(i, j).fill.solid()
                table.cell(i, j).fill.fore_color.rgb = RGBColor(0xD9, 0xE2, 0xF3)

    path = output_dir / "tables_05_large.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: images
# ---------------------------------------------------------------------------
def _create_images(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing embedded images."""
    files: list[Path] = []

    # 1. Single centered image
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    img_bytes = _create_test_image(400, 300, "blue")
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(3), Inches(2), Inches(4), Inches(3))

    path = output_dir / "images_01_centered.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Multiple images tiled
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    colors = ["red", "green", "blue", "yellow"]
    for i, color in enumerate(colors):
        img_bytes = _create_test_image(200, 150, color)
        img_stream = io.BytesIO(img_bytes)
        row = i // 2
        col = i % 2
        slide.shapes.add_picture(
            img_stream,
            Inches(1 + col * 4.5),
            Inches(1 + row * 3.5),
            Inches(3.5),
            Inches(2.5),
        )

    path = output_dir / "images_02_tiled.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Image with text overlay
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    img_bytes = _create_test_image(800, 600, "darkgray")
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(
        img_stream, Inches(0.5), Inches(0.5), Inches(9), Inches(6.5)
    )
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Text overlaying image"
    run.font.size = Pt(36)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.bold = True

    path = output_dir / "images_03_with_text.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Small image with border shape
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    # Border shape
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.5),
        Inches(1.5),
        Inches(5.2),
        Inches(4.2),
    )
    border.fill.background()
    border.line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    border.line.width = Pt(3)

    img_bytes = _create_test_image(300, 225, "coral")
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(2.6), Inches(1.6), Inches(5), Inches(4))

    path = output_dir / "images_04_bordered.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Multiple small images in a row
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    row_colors = ["tomato", "gold", "limegreen", "deepskyblue", "orchid"]
    for i, color in enumerate(row_colors):
        img_bytes = _create_test_image(150, 150, color)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(
            img_stream,
            Inches(0.5 + i * 1.9),
            Inches(3),
            Inches(1.5),
            Inches(1.5),
        )

    path = output_dir / "images_05_row.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: gradients
# ---------------------------------------------------------------------------
def _create_gradients(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing gradient fills."""
    files: list[Path] = []

    # 1. Simple two-color gradient
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5.5),
    )
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = RGBColor(0x41, 0x72, 0xC4)
    shape.fill.gradient_stops[0].position = 0.0
    shape.fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.fill.gradient_stops[1].position = 1.0

    path = output_dir / "gradients_01_two_color.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Three-color gradient
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(5.5),
    )
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].color.rgb = RGBColor(0xFF, 0x00, 0x00)
    stops[0].position = 0.0
    stops[1].color.rgb = RGBColor(0x00, 0x00, 0xFF)
    stops[1].position = 1.0
    # Add middle stop
    stop = stops.add()
    stop.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    stop.position = 0.5

    path = output_dir / "gradients_02_three_color.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Multiple gradient shapes
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    gradient_pairs = [
        (RGBColor(0xFF, 0x63, 0x47), RGBColor(0xFF, 0xA5, 0x00)),
        (RGBColor(0x32, 0xCD, 0x32), RGBColor(0x00, 0x64, 0x00)),
        (RGBColor(0x87, 0xCE, 0xEB), RGBColor(0x00, 0x00, 0x8B)),
    ]
    for i, (c1, c2) in enumerate(gradient_pairs):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + i * 3.2),
            Inches(2),
            Inches(2.8),
            Inches(3),
        )
        shape.fill.gradient()
        shape.fill.gradient_stops[0].color.rgb = c1
        shape.fill.gradient_stops[0].position = 0.0
        shape.fill.gradient_stops[1].color.rgb = c2
        shape.fill.gradient_stops[1].position = 1.0

    path = output_dir / "gradients_03_multiple.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Gradient on oval
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2.5),
        Inches(1.5),
        Inches(5),
        Inches(4.5),
    )
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = RGBColor(0xFF, 0xD7, 0x00)
    shape.fill.gradient_stops[0].position = 0.0
    shape.fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0x45, 0x00)
    shape.fill.gradient_stops[1].position = 1.0

    path = output_dir / "gradients_04_oval.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Dark gradient background with text
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0),
        Emu(0),
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    )
    bg.fill.gradient()
    bg.fill.gradient_stops[0].color.rgb = RGBColor(0x0D, 0x0D, 0x2B)
    bg.fill.gradient_stops[0].position = 0.0
    bg.fill.gradient_stops[1].color.rgb = RGBColor(0x1A, 0x1A, 0x4E)
    bg.fill.gradient_stops[1].position = 1.0
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Gradient Background"
    run.font.size = Pt(40)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.bold = True

    path = output_dir / "gradients_05_dark_bg.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: groups
# ---------------------------------------------------------------------------
def _create_groups(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing grouped shapes and rotation.

    Note: python-pptx has limited group shape support, so we create
    visually grouped shapes positioned together and test rotation.
    """
    files: list[Path] = []

    # 1. Overlapping shapes (visual group)
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    # Background circle
    s1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(3),
        Inches(2),
        Inches(4),
        Inches(4),
    )
    s1.fill.solid()
    s1.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)

    # Foreground rectangle
    s2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4),
        Inches(3),
        Inches(2),
        Inches(2),
    )
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Small accent circle
    s3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.5),
        Inches(3.5),
        Inches(1),
        Inches(1),
    )
    s3.fill.solid()
    s3.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)

    path = output_dir / "groups_01_overlapping.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Rotated shapes
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    rotations = [0, 15, 30, 45, 60, 90]
    for i, rot in enumerate(rotations):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5 + i * 1.5),
            Inches(2.5),
            Inches(1.2),
            Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)
        shape.rotation = rot
        tf = shape.text_frame
        tf.paragraphs[0].text = f"{rot} deg"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    path = output_dir / "groups_02_rotated.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Concentric shapes
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    sizes = [(6, 6), (4.5, 4.5), (3, 3), (1.5, 1.5)]
    shape_colors = [
        RGBColor(0xE8, 0xE8, 0xE8),
        RGBColor(0x5B, 0x9B, 0xD5),
        RGBColor(0x41, 0x72, 0xC4),
        RGBColor(0x2E, 0x53, 0x8C),
    ]
    for (w, h), color in zip(sizes, shape_colors):
        left = Inches((10 - w) / 2)
        top = Inches((7.5 - h) / 2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left,
            top,
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    path = output_dir / "groups_03_concentric.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Grid of small shapes
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    grid_shapes = [MSO_SHAPE.RECTANGLE, MSO_SHAPE.OVAL, MSO_SHAPE.DIAMOND]
    for row in range(3):
        for col in range(5):
            shape_type = grid_shapes[row]
            shape = slide.shapes.add_shape(
                shape_type,
                Inches(0.8 + col * 1.8),
                Inches(1 + row * 2.2),
                Inches(1.2),
                Inches(1.2),
            )
            shape.fill.solid()
            r = (row * 80 + col * 30) % 256
            g = (col * 50 + 100) % 256
            b = (row * 100 + 50) % 256
            shape.fill.fore_color.rgb = RGBColor(r, g, b)

    path = output_dir / "groups_04_grid.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Shapes at various z-orders
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    for i in range(5):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1 + i * 0.8),
            Inches(1 + i * 0.8),
            Inches(4),
            Inches(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            min(255, 65 + i * 40),
            min(255, 114 + i * 20),
            min(255, 196 - i * 20),
        )
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.paragraphs[0].text = f"Layer {i + 1}"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    path = output_dir / "groups_05_z_order.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: layouts
# ---------------------------------------------------------------------------
def _create_layouts(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing master/layout inheritance."""
    files: list[Path] = []

    # 1. Title slide layout
    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.placeholders[0].text = "Presentation Title"
    slide.placeholders[1].text = "Subtitle goes here"

    path = output_dir / "layouts_01_title.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Title + Content layout
    prs = _new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.placeholders[0].text = "Slide Heading"
    tf = slide.placeholders[1].text_frame
    tf.text = "First bullet point"
    tf.add_paragraph().text = "Second bullet point"
    tf.add_paragraph().text = "Third bullet point"

    path = output_dir / "layouts_02_title_content.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Two Content layout
    prs = _new_presentation()
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
        if 0 in slide.placeholders:
            slide.placeholders[0].text = "Two Column Layout"
        if 1 in slide.placeholders:
            slide.placeholders[1].text_frame.text = "Left column content"
        if 2 in slide.placeholders:
            slide.placeholders[2].text_frame.text = "Right column content"
    except (IndexError, KeyError):
        # Fallback: use blank slide with manual two-column
        slide = _add_blank_slide(prs)
        for i, (x, text) in enumerate([(1, "Left column"), (5.5, "Right column")]):
            txBox = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(4), Inches(4))
            txBox.text_frame.text = text

    path = output_dir / "layouts_03_two_content.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Section header layout
    prs = _new_presentation()
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header
        if 0 in slide.placeholders:
            slide.placeholders[0].text = "Section Title"
        if 1 in slide.placeholders:
            slide.placeholders[1].text = "Section description text"
    except (IndexError, KeyError):
        slide = _add_blank_slide(prs)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        txBox.text_frame.text = "Section Title"

    path = output_dir / "layouts_04_section.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Multiple slides with different layouts
    prs = _new_presentation()
    # Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.placeholders[0].text = "Multi-Layout Test"
    slide1.placeholders[1].text = "Testing layout inheritance"

    # Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.placeholders[0].text = "Content Slide"
    slide2.placeholders[1].text_frame.text = "Bullet content"

    # Blank slide with custom content
    slide3 = _add_blank_slide(prs)
    shape = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2),
        Inches(2),
        Inches(6),
        Inches(3.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)
    shape.text_frame.paragraphs[0].text = "Custom blank slide"

    path = output_dir / "layouts_05_multi.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: bullets
# ---------------------------------------------------------------------------
def _create_bullets(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files testing bullet and numbered lists."""
    files: list[Path] = []

    # 1. Simple bullet list
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = [
        "First bullet point with some text content",
        "Second bullet point covering another topic",
        "Third point with technical details included",
        "Fourth point about implementation approach",
        "Fifth and final point wrapping up the list",
    ]
    for text in bullets:
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        p.level = 0

    path = output_dir / "bullets_01_simple.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Nested bullet levels
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    items = [
        ("Main topic one", 0),
        ("Sub-item one-A", 1),
        ("Sub-item one-B", 1),
        ("Detail under one-B", 2),
        ("Main topic two", 0),
        ("Sub-item two-A", 1),
        ("Sub-sub-item", 2),
        ("Deep nested item", 3),
        ("Main topic three", 0),
    ]
    for text, level in items:
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(18 - level * 2)

    path = output_dir / "bullets_02_nested.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Bullet list with bold text
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i in range(6):
        p = tf.add_paragraph()
        # Bold label + normal description
        run_bold = p.add_run()
        run_bold.text = f"Feature {i + 1}: "
        run_bold.font.bold = True
        run_bold.font.size = Pt(18)

        run_normal = p.add_run()
        run_normal.text = f"Description of feature number {i + 1} with details"
        run_normal.font.size = Pt(18)

    path = output_dir / "bullets_03_bold_label.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Colored bullets
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    bullet_colors = [
        RGBColor(0xFF, 0x00, 0x00),
        RGBColor(0xFF, 0xA5, 0x00),
        RGBColor(0xFF, 0xFF, 0x00),
        RGBColor(0x00, 0x80, 0x00),
        RGBColor(0x00, 0x00, 0xFF),
        RGBColor(0x80, 0x00, 0x80),
    ]
    for i, color in enumerate(bullet_colors):
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"Colored bullet point {i + 1}"
        run.font.color.rgb = color
        run.font.size = Pt(22)

    path = output_dir / "bullets_04_colored.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Mixed list with different spacing
    prs = _new_presentation()
    slide = _add_blank_slide(prs)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(6.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    sections = [
        ("Overview", 0, True),
        ("The project aims to convert PPTX to HTML", 1, False),
        ("Using pure Rust for performance", 1, False),
        ("Features", 0, True),
        ("Theme color resolution", 1, False),
        ("Shape rendering", 1, False),
        ("Table support", 1, False),
        ("Status", 0, True),
        ("Currently at version 0.5.0", 1, False),
        ("All major features implemented", 1, False),
    ]
    for text, level, is_header in sections:
        p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(12) if is_header else Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(22) if is_header else Pt(16)
        run.font.bold = is_header

    path = output_dir / "bullets_05_mixed.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Category: mixed (complex slides)
# ---------------------------------------------------------------------------
def _create_mixed(output_dir: Path) -> list[Path]:
    """Generate 5 PPTX files with complex mixed-element slides."""
    files: list[Path] = []

    # 1. Dashboard-style slide
    prs = _new_presentation()
    slide = _add_blank_slide(prs)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Project Dashboard"
    run.font.size = Pt(32)
    run.font.bold = True

    # Metric boxes
    metrics = [("Revenue", "$1.2M"), ("Users", "45.2K"), ("Growth", "+23%")]
    for i, (label, value) in enumerate(metrics):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i * 3.2),
            Inches(1.5),
            Inches(2.8),
            Inches(1.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)

        tf = shape.text_frame
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        r1.font.size = Pt(28)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xD9, 0xE2, 0xF3)

    # Table below
    table_shape = slide.shapes.add_table(
        4, 3, Inches(0.5), Inches(3.5), Inches(9), Inches(3.5)
    )
    table = table_shape.table
    for j, h in enumerate(["Quarter", "Revenue", "Target"]):
        table.cell(0, j).text = h
    data = [
        ["Q1", "$280K", "$250K"],
        ["Q2", "$320K", "$300K"],
        ["Q3", "$340K", "$350K"],
    ]
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            table.cell(i + 1, j).text = val

    path = output_dir / "mixed_01_dashboard.pptx"
    prs.save(str(path))
    files.append(path)

    # 2. Text + Image + Shape composition
    prs = _new_presentation()
    slide = _add_blank_slide(prs)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    run = txBox.text_frame.paragraphs[0].add_run()
    run.text = "Product Feature Highlight"
    run.font.size = Pt(28)
    run.font.bold = True

    # Image on left
    img_bytes = _create_test_image(300, 300, "steelblue")
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), Inches(4), Inches(4))

    # Text on right
    txBox = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    features = [
        "High-performance Rust engine",
        "Accurate color rendering",
        "Complete table support",
        "Theme inheritance chain",
    ]
    for feat in features:
        p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = feat
        run.font.size = Pt(18)

    # Accent shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(6),
        Inches(9),
        Inches(0.1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)

    path = output_dir / "mixed_02_feature.pptx"
    prs.save(str(path))
    files.append(path)

    # 3. Comparison slide
    prs = _new_presentation()
    slide = _add_blank_slide(prs)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Before vs After"
    run.font.size = Pt(32)
    run.font.bold = True

    # Two columns with shapes
    for col, (label, color) in enumerate(
        [
            ("Before", RGBColor(0xCC, 0x33, 0x33)),
            ("After", RGBColor(0x33, 0x99, 0x33)),
        ]
    ):
        x = Inches(0.5 + col * 5)
        # Header shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x,
            Inches(1.5),
            Inches(4.5),
            Inches(1),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = label
        r.font.size = Pt(24)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.bold = True

        # Content below
        txBox = slide.shapes.add_textbox(x, Inches(3), Inches(4.5), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        items = [f"{label} item {i + 1}" for i in range(4)]
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)

    path = output_dir / "mixed_03_comparison.pptx"
    prs.save(str(path))
    files.append(path)

    # 4. Agenda slide with icons
    prs = _new_presentation()
    slide = _add_blank_slide(prs)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Today's Agenda"
    run.font.size = Pt(36)
    run.font.bold = True

    agenda_items = [
        ("Introduction", RGBColor(0x41, 0x72, 0xC4)),
        ("Technical Overview", RGBColor(0xED, 0x7D, 0x31)),
        ("Demo", RGBColor(0x70, 0xAD, 0x47)),
        ("Q&A", RGBColor(0xFB, 0xC0, 0x2D)),
    ]
    for i, (text, color) in enumerate(agenda_items):
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1),
            Inches(1.8 + i * 1.3),
            Inches(0.8),
            Inches(0.8),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        tf = circle.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = str(i + 1)
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Label
        txBox = slide.shapes.add_textbox(
            Inches(2.2),
            Inches(1.9 + i * 1.3),
            Inches(6),
            Inches(0.7),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(22)

    path = output_dir / "mixed_04_agenda.pptx"
    prs.save(str(path))
    files.append(path)

    # 5. Complex multi-element slide
    prs = _new_presentation()
    slide = _add_blank_slide(prs)

    # Background gradient
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0),
        Emu(0),
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    )
    bg.fill.gradient()
    bg.fill.gradient_stops[0].color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    bg.fill.gradient_stops[0].position = 0.0
    bg.fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.fill.gradient_stops[1].position = 1.0
    bg.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0),
        Emu(0),
        SLIDE_WIDTH,
        Inches(1.2),
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(0x2E, 0x3A, 0x4F)
    title_bar.line.fill.background()
    tf = title_bar.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = "Architecture Overview"
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.bold = True

    # Process flow boxes
    steps = ["Parse", "Model", "Resolve", "Render"]
    step_colors = [
        RGBColor(0x41, 0x72, 0xC4),
        RGBColor(0xED, 0x7D, 0x31),
        RGBColor(0x70, 0xAD, 0x47),
        RGBColor(0xFB, 0xC0, 0x2D),
    ]
    for i, (step, color) in enumerate(zip(steps, step_colors)):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5 + i * 2.4),
            Inches(2),
            Inches(2),
            Inches(1.5),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = step
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Arrow between boxes (except last)
        if i < 3:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(2.5 + i * 2.4),
                Inches(2.5),
                Inches(0.4),
                Inches(0.5),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(0x88, 0x88, 0x88)

    # Description text at bottom
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "The pptx2html-rs pipeline processes PPTX files through four stages: "
    run.font.size = Pt(14)
    run2 = p.add_run()
    run2.text = (
        "SAX parsing, model construction, property resolution, and HTML rendering."
    )
    run2.font.size = Pt(14)
    run2.font.italic = True

    path = output_dir / "mixed_05_architecture.pptx"
    prs.save(str(path))
    files.append(path)

    return files


# ---------------------------------------------------------------------------
# Main orchestrator
# ---------------------------------------------------------------------------
CATEGORY_GENERATORS = {
    "basic_text": _create_basic_text,
    "shapes": _create_shapes,
    "theme_colors": _create_theme_colors,
    "tables": _create_tables,
    "images": _create_images,
    "gradients": _create_gradients,
    "groups": _create_groups,
    "layouts": _create_layouts,
    "bullets": _create_bullets,
    "mixed": _create_mixed,
}


def create_golden_set(
    output_dir: Path, categories: list[str] | None = None
) -> dict[str, list[Path]]:
    """Generate all golden PPTX test files.

    Args:
        output_dir: Directory to save generated PPTX files.
        categories: If specified, only generate these categories.

    Returns:
        Dict mapping category name -> list of generated file paths.
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    generators = CATEGORY_GENERATORS
    if categories:
        generators = {k: v for k, v in generators.items() if k in categories}

    results: dict[str, list[Path]] = {}
    total = 0

    for category, gen_fn in generators.items():
        logger.info("Generating category: %s", category)
        try:
            paths = gen_fn(output_dir)
            results[category] = paths
            total += len(paths)
            logger.info("  -> %d files generated", len(paths))
        except Exception as exc:
            logger.error("Failed to generate %s: %s", category, exc)

    logger.info("Total: %d golden PPTX files generated in %s", total, output_dir)
    return results


def main() -> None:
    """CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Generate golden PPTX test set for evaluation"
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Output directory (default: evaluate/golden_set/)",
    )
    parser.add_argument(
        "--categories",
        nargs="*",
        choices=list(CATEGORY_GENERATORS.keys()),
        help="Generate only specified categories",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Enable debug logging",
    )

    args = parser.parse_args()

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s: %(message)s",
    )

    output_dir = args.output
    if output_dir is None:
        output_dir = Path(__file__).resolve().parent / "golden_set"

    create_golden_set(output_dir=output_dir, categories=args.categories)


if __name__ == "__main__":
    main()
